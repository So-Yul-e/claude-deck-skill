# -*- coding: utf-8 -*-
"""
deck 스킬 — 고정 테마 빌더.

색·여백·타이포는 여기 잠겨 있다. 덱 작성자는 빌더만 호출한다.
시각 빌더(flow·timeline·chart·matrix·shots·progress)는 전부 **도형·네이티브
차트**로 그린다. 래스터 이미지가 아니라서 PowerPoint에서 그대로 편집되고,
테마 색을 따라간다.

    import sys; sys.path.insert(0, "<skill_dir>")
    from deck import Deck
    d = Deck(palette="indigo", footer="프로젝트 · 팀")
    d.cover("제목", "부제")
    d.flow("유저 플로우", [("발견","웹→설치"), ("온보딩","60초"), ("추천","10초")])
    d.save("out.pptx")
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── 팔레트 ────────────────────────────────────────────────────────────
PALETTES = {
    "indigo": dict(deep="141A5C", key="3D52D5", sub="7B68C8", tint="EFF1FC",
                   pale="CED5F8", rule="E5E5EA"),
    "navy":   dict(deep="0F1E38", key="1F4E9C", sub="3A6FC4", tint="EDF2F9",
                   pale="C6D6EA", rule="E5E5EA"),
    "mono":   dict(deep="1A1A1A", key="2C2C2C", sub="6E6E73", tint="F2F2F4",
                   pale="D2D2D7", rule="E5E5EA"),
}
INK = RGBColor(0x1C, 0x1C, 0x1E)
MUTE = RGBColor(0x6E, 0x6E, 0x73)   # 흰 배경 5.07:1 (구 8E8E93은 3.26:1로 본문 기준 미달)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
FONT_DIR = SKILL_DIR / "assets" / "fonts"
FONT_CANDIDATES = {
    "Pretendard": [
        FONT_DIR / "Pretendard-Regular.otf",
        Path.home() / "Library" / "Fonts" / "Pretendard-Regular.otf",
        Path("/Library/Fonts/Pretendard-Regular.otf"),
        Path.home() / "AppData" / "Local" / "Microsoft" / "Windows" / "Fonts" / "Pretendard-Regular.otf",
    ],
}


def _font():
    for path in FONT_CANDIDATES["Pretendard"]:
        if path.exists():
            return "Pretendard"
    raise RuntimeError(
        "Pretendard가 필요합니다. 먼저 deck 스킬의 "
        "'bash deck/scripts/deps-macos.sh --check' 또는 "
        "'powershell -File deck/scripts/deps-windows.ps1 -Check'를 실행하세요."
    )


FONT = _font()

# ── 텍스트 실측 ───────────────────────────────────────────────────────
# 칩 폭·자동 축소를 "글자수 × 상수"로 근사하면 한글/영문/숫자 혼용에서 반드시
# 틀린다. Pillow로 실제 폰트를 재고, 폰트를 못 열면 보수적 근사로 폴백한다.
_FONT_FILES = {
    "Pretendard": {
        "regular": next((str(p) for p in FONT_CANDIDATES["Pretendard"] if p.exists()), None),
        "bold": str(FONT_DIR / "Pretendard-Bold.otf") if (FONT_DIR / "Pretendard-Bold.otf").exists() else None,
    },
}
_MEASURE_CACHE = {}


def measure_pt(text, size_pt, bold=False):
    """텍스트 렌더 폭을 pt로 반환."""
    key = (text, size_pt, bold)
    if key in _MEASURE_CACHE:
        return _MEASURE_CACHE[key]
    val = None
    paths = _FONT_FILES.get(FONT, {})
    path = paths.get("bold" if bold and paths.get("bold") else "regular")
    if path and os.path.exists(path):
        try:
            from PIL import ImageFont
            idx = 2 if (bold and path.endswith(".ttc")) else 0
            f = ImageFont.truetype(path, int(size_pt * 4), index=idx)
            val = f.getlength(str(text)) / 4.0
        except Exception:
            val = None
    if val is None:
        # 폴백: 한글 1.0em · 그 외 0.55em
        w = sum(1.0 if ord(c) > 0x2000 else 0.55 for c in str(text))
        val = w * size_pt
    if bold:
        val *= 1.04
    _MEASURE_CACHE[key] = val
    return val


def measure_emu(text, size_pt, bold=False):
    return Emu(int(measure_pt(text, size_pt, bold) * 12700))


# ── 폰트 안전 문자 ───────────────────────────────────────────────────
# 폰트나 렌더러가 지원하지 않는 글리프는 대체 폰트로 바뀌어 자간·굵기·높이가
# 어긋난다. 특히 짝을 이루는 기호(▓/░ 진행률)에서 한쪽만 폴백되면 눈에 띈다.
# 아래는 Pretendard + Office/LibreOffice 조합에서 실측한 안전 목록이다.
SAFE_GLYPHS = {
    "채움": "█■●━", "빈칸": "░□○─·", "순환": "⇄→⤴",
    "닫기": "×✖X", "경고": "!△▲※", "화살": "→←↑↓",
}
UNSAFE_GLYPHS = "▓▪▬↺↻⟲♺✕╳⚠▫"   # 폴백 발생 — 쓰지 말 것


def check_glyphs(text):
    """폰트 폴백을 유발하는 문자를 돌려준다. 덱 작성 후 자가 점검용."""
    return sorted({c for c in str(text) if c in UNSAFE_GLYPHS})

# 타이포 — 인접 단계 1.35배 이상 (SKILL.md §1)
T_COVER, T_SECNUM = 50, 68
T_TITLE = T_SEC = 30          # 동시 등장 없음 — 상수를 둘로 둘 이유가 없다
T_LEAD = T_BODY = 19          # 리드/본문은 크기가 아니라 색(MUTE/INK)으로 가른다
T_SMALL, T_META = 14, 11      # 19/14 = 1.36 · 14/11 = 1.27(둘 다 캡션 급)
T_VALUE = 38                  # cards 값 전용 — 종전 30 하드코딩은 T_TITLE과 동급이었다

SW, SH = Inches(13.333), Inches(7.5)
M = Inches(0.85)
CW = SW - M * 2
BOTTOM = SH - Inches(1.15)   # 푸터 룰 위. 콘텐츠는 여기까지 쓴다.


class Deck:
    def __init__(self, palette="indigo", footer=""):
        p = PALETTES.get(palette, palette) if isinstance(palette, str) else palette
        self.C = {k: RGBColor.from_string(v) for k, v in p.items()}
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SW, SH
        self.footer = footer
        self._n = 0

    # ── 내부 ──────────────────────────────────────────────────────
    def _rect(self, s, x, y, w, h, fill=None, line=None, radius=False):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        sp = s.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
        if radius:
            try:
                sp.adjustments[0] = 0.08
            except Exception:
                pass
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(1)
        sp.shadow.inherit = False
        return sp

    def _oval(self, s, x, y, w, h, fill):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
        sp.line.fill.background(); sp.shadow.inherit = False
        return sp

    def _txt(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, ls=1.35):
        box = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        i = 0
        for run in runs:
            t, sz, b, c = run[0], run[1], run[2], run[3]
            url = run[4] if len(run) > 4 else None
            # 문자열 안의 \n은 문단으로 쪼갠다. python-pptx는 개행을 <a:br/>로
            # 바꾸지 않고 <a:t>에 그대로 넣기 때문에, 렌더러가 줄은 바꾸면서도
            # 첫 줄에는 문단 정렬을 적용하지 않는다. 타임라인 라벨 1행만
            # 노드에서 44~78pt 왼쪽으로 밀려 계단처럼 보이던 원인이 이것이다.
            for line in str(t).split("\n"):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = align; p.line_spacing = ls
                r = p.add_run(); r.text = line
                r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = b
                r.font.color.rgb = c
                if url:
                    r.hyperlink.address = url
                    r.font.underline = True
                i += 1
        return box

    def _blank(self, dark=False):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._last = s
        if dark:
            self._rect(s, 0, 0, SW, SH, fill=self.C["deep"])
        return s

    def _head(self, s, title, eyebrow=None, lead=None):
        """본문 슬라이드 공통 헤더. 반환값 = 콘텐츠 시작 y."""
        self._n += 1
        if eyebrow:
            self._rect(s, M, Inches(0.72), Inches(0.055), Inches(0.42), fill=self.C["key"])
            self._txt(s, M + Inches(0.22), Inches(0.72), Inches(7), Inches(0.42),
                      [(eyebrow.upper(), T_META, True, self.C["key"])],
                      anchor=MSO_ANCHOR.MIDDLE)
        self._txt(s, M, Inches(1.26), CW, Inches(0.8), [(title, T_TITLE, True, INK)])
        y = Inches(2.25)
        if lead:
            self._txt(s, M, Inches(2.05), CW, Inches(0.5), [(lead, T_LEAD, False, MUTE)])
            y = Inches(2.75)
        self._rect(s, M, SH - Inches(0.92), CW, Emu(9525), fill=self.C["rule"])
        self._txt(s, M, SH - Inches(0.78), Inches(9), Inches(0.3),
                  [(self.footer, T_META, False, MUTE)])
        self._txt(s, SW - M - Inches(1.2), SH - Inches(0.78), Inches(1.2), Inches(0.3),
                  [(f"{self._n:02d}", T_META, True, self.C["key"])], align=PP_ALIGN.RIGHT)
        return y

    def _accent(self, i, hot=None):
        """기본은 전 항목 key. hot을 준 항목만 key, 나머지는 sub.

        종전에는 i==0만 key였는데, 인덱스 0에 의미가 있는 경우가 없어서
        색이 "첫 항목이 특별하다"는 거짓 정보를 만들었다(전부 100%인 progress에서
        첫 행만 색이 다른 식).
        """
        if hot is None:
            return self.C["key"]
        return self.C["key"] if i == hot else self.C["sub"]

    # ── 텍스트 슬라이드 ────────────────────────────────────────────
    def cover(self, title, subtitle="", meta=""):
        s = self._blank(dark=True)
        self._rect(s, M, Inches(2.35), Inches(0.9), Inches(0.06), fill=self.C["sub"])
        self._txt(s, M, Inches(2.72), CW, Inches(1.4),
                  [(title, T_COVER, True, WHITE)], ls=1.15)
        if subtitle:
            self._txt(s, M, Inches(4.22), CW, Inches(0.6),
                      [(subtitle, T_LEAD, False, self.C["pale"])])
        if meta:
            self._txt(s, M, SH - Inches(1.05), CW, Inches(0.4),
                      [(meta, T_SMALL, False, self.C["pale"])])
        return s

    def section(self, number, title, lead=""):
        s = self._blank(dark=True)
        self._txt(s, M, Inches(2.05), Inches(3), Inches(1.2),
                  [(str(number), T_SECNUM, True, self.C["sub"])], ls=1.0)
        self._txt(s, M, Inches(3.3), CW, Inches(0.9), [(title, T_SEC, True, WHITE)])
        if lead:
            self._txt(s, M, Inches(4.3), Inches(8.4), Inches(1.0),
                      [(lead, T_LEAD, False, self.C["pale"])])
        return s

    def statement(self, text, sub=""):
        s = self._blank()
        self._head(s, "")
        # 가용 영역(헤더 아래 ~ 푸터 위) 세로 중앙. 위에 붙이면 아래가 통째로 빈다.
        lines = text.count("\n") + 1
        # 한 줄 높이를 0.62in(44.6pt)로 가정하면 실제 30pt×1.3=39pt와 어긋나
        # 세로선이 글자보다 길어지고, top 정렬 탓에 위로도 튀어나온다.
        # 실제 행 높이로 계산하고 선과 텍스트가 같은 박스를 공유하게 한다.
        line_h = Pt(T_SEC * 1.3)
        block_h = line_h * lines
        block = block_h + (Inches(0.9) if sub else 0)
        top = Inches(1.05) + (BOTTOM - Inches(1.05) - block) / 2
        self._rect(s, M, top, Inches(0.06), block_h, fill=self.C["key"])
        self._txt(s, M + Inches(0.45), top, CW - Inches(0.45), block_h,
                  [(text, T_SEC, True, INK)], ls=1.3, anchor=MSO_ANCHOR.MIDDLE)
        if sub:
            self._txt(s, M + Inches(0.45), top + block_h + Inches(0.24),
                      Inches(10.2), Inches(0.8), [(sub, T_BODY, False, MUTE)])
        return s

    def bullets(self, title, items, eyebrow=None, lead=None, hot=None):
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        row = min(Inches(0.9), (BOTTOM - y) / max(1, len(items)))
        line_h = Inches(0.34)          # T_BODY 19pt 한 줄 높이
        dot = Inches(0.14)
        for i, it in enumerate(items):
            # 마커는 텍스트 첫 줄의 세로 중심에 맞춘다(수동 오프셋 금지).
            self._rect(s, M, y + (line_h - dot) / 2, dot, dot,
                       fill=self._accent(i, hot))
            self._txt(s, M + Inches(0.4), y, CW - Inches(0.4), line_h,
                      [(it, T_BODY, False, INK)], anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
            y += row
        return s

    def cards(self, title, items, eyebrow=None, hot=None, lead=None):
        """items = [(라벨, 값, 설명)] 최대 4."""
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        n = max(1, min(len(items), 4))
        g = Inches(0.28); cw = (CW - g * (n - 1)) / n
        # 콘텐츠 실측 높이로 고정하고 남는 세로는 중앙 배치로 흡수한다.
        # 종전 3.2in 고정은 카드마다 하단 1in가 빈 통으로 남았다.
        ch_ = min(Inches(2.55), BOTTOM - y)
        y = y + (BOTTOM - y - ch_) / 2
        # 값 크기는 카드 폭과 글자수로 정한다. 고정하면 "500~1,000" 같은 값이
        # 두 줄로 깨져 아래 설명과 겹친다(실측 결함).
        # 실제 텍스트 박스 폭(cw-0.6in)을 pt로 환산하고, 한글·기호 혼용을 고려해
        # 글자당 0.62em으로 잡는다. 종전 1.75 계수는 경계에서 한 글자가 넘쳐
        # "500~1,000"이 두 줄로 깨졌다(실측).
        box_pt = (cw - Inches(0.44)) / Emu(12700)
        vsize = T_VALUE
        widest = max(items[:n], key=lambda t: measure_pt(t[1], T_VALUE, True))[1]
        while vsize > 20 and measure_pt(widest, vsize, True) > box_pt:
            vsize -= 1
        for i, (lab, val, desc) in enumerate(items[:n]):
            x = M + (cw + g) * i
            self._rect(s, x, y, cw, ch_, fill=self.C["tint"])
            self._rect(s, x, y, cw, Inches(0.055), fill=self._accent(i, hot))
            self._txt(s, x + Inches(0.22), y + Inches(0.3), cw - Inches(0.44), Inches(0.34),
                      [(lab, T_META, True, self.C["key"])])
            self._txt(s, x + Inches(0.22), y + Inches(0.68), cw - Inches(0.44), Inches(0.8),
                      [(val, vsize, True, self.C["deep"])], ls=1.1)
            self._txt(s, x + Inches(0.22), y + Inches(1.52), cw - Inches(0.44), Inches(0.7),
                      [(desc, T_SMALL, False, MUTE)])
        return s

    def table(self, title, headers, rows, eyebrow=None, col_ratio=None, lead=None):
        """행·열 대조표.

        python-pptx 기본값은 셀 텍스트가 **좌측 상단에 붙고 안쪽 여백이 0**이라
        그대로 두면 표가 싸구려로 보인다. 세로 중앙 정렬 + 여백을 명시하고,
        행 높이를 가용 공간에 분배해 아래가 비지 않게 한다.
        """
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        nr, nc = len(rows) + 1, len(headers)
        avail = BOTTOM - y
        # 헤더는 조금 낮게, 본문 행은 남는 높이를 균등 분배(하한 0.5in)
        hdr_h = Inches(0.52)
        body_h = min(Inches(0.86), max(Inches(0.46), (avail - hdr_h) / max(1, len(rows))))
        total = hdr_h + body_h * len(rows)
        tbl = s.shapes.add_table(nr, nc, int(M), int(y), int(CW), int(total)).table
        tbl.first_row = True
        tbl.rows[0].height = int(hdr_h)
        for r_ in range(1, nr):
            tbl.rows[r_].height = int(body_h)
        if col_ratio:
            tot = sum(col_ratio)
            for i, r in enumerate(col_ratio):
                tbl.columns[i].width = Emu(int(CW * r / tot))

        def cell(rr, cc, text, bold, bg, fg, size):
            c = tbl.cell(rr, cc)
            c.text = str(text)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            c.vertical_anchor = MSO_ANCHOR.MIDDLE          # 위로 붙는 것 방지
            c.margin_left = c.margin_right = Emu(109728)   # 0.12in
            c.margin_top = c.margin_bottom = Emu(73152)    # 0.08in
            p_ = c.text_frame.paragraphs[0]
            p_.line_spacing = 1.3
            p_.runs[0].font.name = FONT
            p_.runs[0].font.size = Pt(size)
            p_.runs[0].font.bold = bold
            p_.runs[0].font.color.rgb = fg

        for c_, htxt in enumerate(headers):
            cell(0, c_, htxt, True, WHITE, self.C["key"], T_META)
        self._rect(s, M, y + hdr_h - Emu(19050), CW, Emu(19050), fill=self.C["key"])
        for r_, row in enumerate(rows, 1):
            for c_, v in enumerate(row):
                cell(r_, c_, v, False, WHITE if r_ % 2 else self.C["tint"], INK, T_SMALL)
        return s

    def deflist(self, title, items, eyebrow=None, lead=None):
        """라벨 정의 목록. items = [(라벨, 서술)] 또는 [(라벨, 서술, url)] 5~6행 상한.

        url을 주면 라벨이 클릭 가능한 하이퍼링크가 된다. 참고자료 슬라이드에서
        경로만 적어두면 아무도 못 찾는다. 링크를 건다.

        **2열 표를 쓰지 말 것.** 열이 하나면 대조할 게 없어서 표가 아니라 헤더가
        붙은 불릿이다. 헤더 행은 정보가 0인데 시각 무게만 먹고, 교대 배경까지
        붙으면 밀도가 과장된다. 여기는 헤더도 교대 배경도 없다.
        """
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        n = max(1, len(items))
        row = (BOTTOM - y) / n
        lab_w = CW * 0.28
        # 라벨이 폭을 넘으면 두 줄이 되어 다음 행과 겹친다. 실측으로 자동 축소.
        lab_box = (lab_w - Inches(0.2)) / Emu(12700)
        lsize = T_BODY
        widest = max((it[0] for it in items), key=lambda t: measure_pt(t, T_BODY, True))
        while lsize > 13 and measure_pt(widest, lsize, True) > lab_box:
            lsize -= 1
        for i, item in enumerate(items):
            lab, body = item[0], item[1]
            url = item[2] if len(item) > 2 else None
            ry = y + row * i
            if i:
                self._rect(s, M, ry, CW, Emu(9525), fill=self.C["rule"])
            # 라벨과 서술은 글자 크기·행간이 달라서 top 정렬하면 첫 줄 baseline이
            # 어긋난다(16pt×1.35 vs 13pt×1.45). 같은 행 높이 안에서 세로 중앙에
            # 맞춰야 크기가 달라도 눈에는 한 줄로 보인다.
            self._txt(s, M, ry, lab_w - Inches(0.2), row,
                      [(lab, lsize, True, self.C["key"]) if not url
                       else (lab, lsize, True, self.C["key"], url)],
                      anchor=MSO_ANCHOR.MIDDLE)
            self._txt(s, M + lab_w, ry, CW - lab_w, row,
                      [(body, T_SMALL, False, INK)], ls=1.45,
                      anchor=MSO_ANCHOR.MIDDLE)
        return s

    def tree(self, title, groups, eyebrow=None, lead=None, orientation="row", hot=None):
        """계층·클러스터. groups = [(부모, [자식, ...], 캡션), ...] 최대 5.

        orientation="stack"이면 위아래로 쌓인 레이어(아키텍처·문제 레이어)가 된다.

        표로 대신하면 부모가 반복 텍스트로 눌려 관계가 지워지고, flow로 대신하면
        없는 순서를 주장하게 된다. IA·기능 클러스터는 순서가 없는 포함 관계다.

        6개 이상은 **자르지 않고 예외를 낸다.** 종전 `min(len, 4)`는 5번째 그룹을
        말없이 버렸고, 제목이 "5대 클러스터"인 슬라이드에 클러스터가 4개만 그려진
        채 배포 직전까지 아무도 몰랐다. 넘치면 조용히 사라지는 것보다 만드는
        쪽이 멈추는 게 낫다.
        """
        if len(groups) > 5:
            raise ValueError(
                f"tree()는 그룹 5개까지다(받은 값 {len(groups)}개). "
                "잘라내면 내용이 소리 없이 사라진다 — 슬라이드를 나누거나 "
                "table()·deflist()로 형태를 바꿀 것.")
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        n = max(1, len(groups))
        row = (BOTTOM - y) / n
        # 5행이면 행 높이가 라벨 블록(제목 0.34 + 캡션 0.26)보다 얇아진다.
        # 칩과 캡션을 함께 줄여 넘침 대신 밀도로 흡수한다.
        tight = row < Inches(0.66)
        lab_w = CW * 0.22
        for i, gitem in enumerate(groups[:n]):
            parent, children = gitem[0], gitem[1]
            # 문자열을 넘기면 파이썬이 글자 단위로 순회해 한 글자짜리 칩이 쏟아진다.
            # 공백 구분 문자열도 받아준다.
            if isinstance(children, str):
                children = children.split()
            cap = gitem[2] if len(gitem) > 2 else ""
            ry = y + row * i
            # tight(5행)면 라벨 블록·칩을 함께 줄인다. 한쪽만 줄이면 액센트 바가
            # 글자 높이와 어긋나 "바가 라벨과 안 맞는" 그 증상이 다시 나온다.
            t_h = Inches(0.30) if tight else Inches(0.34)
            c_off = Inches(0.30) if tight else Inches(0.44)
            c_h = Inches(0.22) if tight else Inches(0.26)
            # 액센트 바는 행 전체가 아니라 라벨 블록(제목 + 캡션) 높이에 맞춘다.
            lab_h = t_h + (c_h if cap else Emu(0))
            self._rect(s, M, ry + Inches(0.06), Inches(0.055), lab_h,
                       fill=self._accent(i, hot))
            self._txt(s, M + Inches(0.22), ry + Inches(0.04), lab_w - Inches(0.3), t_h,
                      [(parent, T_BODY, True, INK)], ls=1.2)
            if cap:
                self._txt(s, M + Inches(0.22), ry + c_off, lab_w - Inches(0.3),
                          c_h + Inches(0.06), [(cap, T_META, False, MUTE)])
            # 자식은 칩으로 wrap 배치
            chip_h = Inches(0.38) if tight else Inches(0.42)
            gap = Inches(0.12) if tight else Inches(0.14)
            cx, cy = M + lab_w, ry + (Inches(0.06) if tight else Inches(0.1))
            for ch in children:
                # 실측 폭 + 좌우 패딩 0.16in씩
                w = measure_emu(ch, T_SMALL) + Inches(0.32)
                if cx + w > M + CW:
                    cx = M + lab_w; cy += chip_h + Inches(0.12)
                self._rect(s, cx, cy, w, chip_h, fill=self.C["tint"], radius=True)
                # 수동 y 오프셋으로 맞추면 폰트·크기가 바뀔 때마다 틀어진다.
                # 칩 전체 높이를 주고 세로 중앙 앵커로 맡긴다.
                self._txt(s, cx, cy, w, chip_h, [(str(ch), T_SMALL, False, INK)],
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
                cx += w + gap
        return s

    # ── 시각 빌더 ─────────────────────────────────────────────────
    def flow(self, title, steps, eyebrow=None, lead=None, per_row=4, loop=False):
        """단계 다이어그램. steps = [(제목, 설명), ...]

        loop=True면 마지막에서 첫 단계로 돌아오는 커넥터를 그린다. 순환 구조를
        직선으로 그리면 "쓸수록 깊어진다"는 주장을 그림이 부정한다.

        카드 제목은 **한글 8자·영문 12자**를 넘기지 말 것. 넘치면 두 줄이 되어
        설명과 겹친다(§2 판정 조건).
        """
        s = self._blank()
        y0 = self._head(s, title, eyebrow, lead)
        n = len(steps)
        rows = (n + per_row - 1) // per_row
        arrow = Inches(0.42); g = Inches(0.1)
        cw = (CW - (arrow + g * 2) * (per_row - 1)) / per_row
        avail = BOTTOM - y0 - (Inches(0.55) if loop else 0)
        vgap = Inches(0.5)
        bh = min(Inches(1.35), (avail - vgap * (rows - 1)) / rows)
        # 한 행뿐이면 남는 세로의 중앙에 놓는다(위에 붙어 아래가 비는 것 방지)
        top = y0 + (avail - bh) / 2 if rows == 1 else y0

        last_xy = None
        for i, (head, desc) in enumerate(steps):
            r, c = divmod(i, per_row)
            x = M + (cw + arrow + g * 2) * c
            y = top + (bh + vgap) * r
            self._rect(s, x, y, cw, bh, fill=self.C["tint"], radius=True)
            # 제목 2줄까지 허용하고 설명을 그 아래로 — 겹침 방지
            self._txt(s, x + Inches(0.14), y + bh / 2 - Inches(0.56), cw - Inches(0.28),
                      Inches(0.62), [(f"{i+1}. {head}", T_SMALL, True, INK)],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, ls=1.15)
            self._txt(s, x + Inches(0.14), y + bh / 2 + Inches(0.02), cw - Inches(0.28),
                      Inches(0.6), [(desc, T_META, False, MUTE)],
                      align=PP_ALIGN.CENTER, ls=1.25)
            if c < per_row - 1 and i < n - 1:
                self._txt(s, x + cw + g, y + bh / 2 - Inches(0.2), arrow, Inches(0.4),
                          [("→", 20, False, self.C["sub"])], align=PP_ALIGN.CENTER)
            if i == n - 1:
                last_xy = (x, y)

        if loop and last_xy:
            lx, ly = last_xy
            by = ly + bh + Inches(0.3)
            self._rect(s, M + cw / 2, by, lx + cw / 2 - (M + cw / 2), Emu(19050),
                       fill=self.C["sub"])
            self._rect(s, M + cw / 2, ly + bh, Emu(19050), by - (ly + bh), fill=self.C["sub"])
            self._rect(s, lx + cw / 2, ly + bh, Emu(19050), by - (ly + bh), fill=self.C["sub"])
            self._txt(s, M + cw / 2 + Inches(0.14), by + Inches(0.06), Inches(3.4), Inches(0.3),
                      [("→  다음 방문에 반영", T_META, True, self.C["sub"])])
        return s

    def timeline(self, title, milestones, eyebrow=None, lead=None):
        """마일스톤 축. milestones = [(날짜, 라벨, 강조bool), ...]"""
        s = self._blank()
        y0 = self._head(s, title, eyebrow, lead)
        # 블록(날짜 0.72 + 축 + 라벨 0.9)을 가용 영역 세로 중앙에 놓는다.
        axis_y = y0 + (BOTTOM - y0) / 2 - Inches(0.1)
        step_ = CW / max(1, len(milestones))
        self._rect(s, M + step_ / 2 - Inches(0.35), axis_y,
                   CW - step_ + Inches(0.70), Emu(19050), fill=self.C["rule"])
        n = max(1, len(milestones))
        step = CW / n
        for i, m in enumerate(milestones):
            date, label = m[0], m[1]
            hot = m[2] if len(m) > 2 else False
            cx = M + step * i + step / 2
            col = self.C["key"] if hot else self.C["sub"]
            d = Inches(0.19)
            self._oval(s, cx - d / 2, axis_y - d / 2 + Emu(9525), d, d, col)
            self._txt(s, cx - step / 2, axis_y - Inches(0.72), step, Inches(0.45),
                      [(date, T_LEAD, True, col)], align=PP_ALIGN.CENTER)
            self._txt(s, cx - step / 2, axis_y + Inches(0.3), step, Inches(1.0),
                      [(label, T_SMALL, hot, INK)], align=PP_ALIGN.CENTER, ls=1.3)
        return s

    def chart(self, title, items, kind="bar", eyebrow=None, lead=None):
        """네이티브 차트. items = [(라벨, 값), ...] · kind = bar|column|donut"""
        s = self._blank()
        y0 = self._head(s, title, eyebrow, lead)
        data = CategoryChartData()
        data.categories = [i[0] for i in items]
        data.add_series("값", tuple(float(i[1]) for i in items))
        t = {"bar": XL_CHART_TYPE.BAR_CLUSTERED,
             "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
             "donut": XL_CHART_TYPE.DOUGHNUT}[kind]
        h = BOTTOM - y0
        # 도넛은 전폭 프레임을 주면 원은 가운데, 범례는 슬라이드 끝으로 밀려
        # 시선이 크게 튄다. 폭을 좁혀 원과 범례를 붙이고 가운데 정렬한다.
        cw_ = CW * 0.66 if kind == "donut" else CW
        cx_ = M + (CW - cw_) / 2
        gf = s.shapes.add_chart(t, int(cx_), int(y0), int(cw_), int(h), data)
        ch = gf.chart
        ch.font.name = FONT; ch.font.size = Pt(T_SMALL); ch.font.color.rgb = INK
        ch.has_title = False   # 시리즈명이 제목으로 떠서 지운다(슬라이드 타이틀과 중복)
        # python-pptx 네이티브 차트는 Office 기본 회색 그리드라인·축선을 그대로
        # 그린다. §1 금지 목록의 "기본 Office 테마 색"에 정확히 해당한다.
        if kind != "donut":
            va = ch.value_axis
            va.has_major_gridlines = False
            va.visible = False
            ca = ch.category_axis
            ca.format.line.color.rgb = self.C["rule"]
            ca.has_major_gridlines = False
        plot = ch.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.name = FONT
        dl.font.size = Pt(T_SMALL)
        dl.font.bold = True
        # 도넛 라벨은 진한 조각 위에 얹히므로 어두운 글씨는 읽히지 않는다.
        dl.font.color.rgb = WHITE if kind == "donut" else self.C["deep"]
        if kind == "donut":
            ch.has_legend = True
            # 우측 범례는 원에서 멀어진다. 아래에 두면 시선 이동이 짧고 대칭이다.
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
            # 범례 폰트를 명시하지 않으면 렌더러가 숫자만 다른 폰트로 폴백해
            # 위첨자처럼 작게 그린다(LibreOffice 실측).
            lf = ch.legend.font
            lf.name = FONT; lf.size = Pt(T_SMALL); lf.bold = False; lf.color.rgb = INK
            pts = ch.plots[0].series[0].points
            for i, pt in enumerate(pts):
                pt.format.fill.solid()
                base = [self.C["key"], self.C["sub"], self.C["pale"], self.C["rule"]]
                pt.format.fill.fore_color.rgb = base[i % len(base)]
        else:
            ch.has_legend = False
            sr = ch.plots[0].series[0]
            sr.format.fill.solid(); sr.format.fill.fore_color.rgb = self.C["key"]
            ch.plots[0].gap_width = 60
        return s

    def progress(self, title, items, eyebrow=None, lead=None, hot=None):
        """진척 바. items = [(라벨, 0~100, 보조텍스트), ...]"""
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        row = min(Inches(0.95), (BOTTOM - y) / max(1, len(items)))
        for i, it in enumerate(items):
            lab, pct = it[0], max(0, min(100, float(it[1])))
            note = it[2] if len(it) > 2 else ""
            band = Inches(0.36)        # 라벨·바·노트가 공유하는 밴드
            bar_h = Inches(0.22)
            self._txt(s, M, y, Inches(3.2), band,
                      [(lab, T_BODY, True, INK)], anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
            bx = M + Inches(3.4); bw = CW - Inches(3.4) - Inches(1.5)
            by_ = y + (band - bar_h) / 2      # 밴드 세로 중앙
            self._rect(s, bx, by_, bw, bar_h, fill=self.C["tint"])
            self._rect(s, bx, by_, bw * pct / 100, bar_h, fill=self._accent(i, hot))
            self._txt(s, SW - M - Inches(1.4), y, Inches(1.4), band,
                      [(note or f"{pct:.0f}%", T_SMALL, True, self.C["key"])],
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
            y += row
        return s

    def matrix(self, title, x_label, y_label, items, eyebrow=None, quadrants=None, lead=None):
        """2x2 포지셔닝 맵. items = [(이름, x 0~1, y 0~1, 강조bool), ...]

        배치 규칙 — 좌우 flip을 쓰지 않는다.
        정사각형을 고집하면 세로 제약 때문에 박스가 작아지고, 그러면 라벨이
        박스를 넘쳐 대부분 왼쪽으로 뒤집히면서 점과 글자의 좌우 관계가 뒤죽박죽이
        된다. 포지셔닝 맵은 정사각형일 이유가 없으므로 **콘텐츠 폭을 다 쓰는
        직사각형**으로 두고, 라벨은 **점 아래 중앙**에 놓는다. 좌우로 넘칠 일이
        없어지고 겹침은 세로 오프셋만으로 해결된다.
        """
        s = self._blank()
        y0 = self._head(s, title, eyebrow)
        LBL = Inches(1.75)                     # y축 라벨 폭
        bx, by = M + LBL, y0
        bw, bh = CW - LBL, BOTTOM - y0 - Inches(0.42)   # 0.42 = x축 라벨 자리

        self._rect(s, bx, by, bw, bh, fill=self.C["tint"])
        self._rect(s, bx + bw / 2, by, Emu(9525), bh, fill=self.C["rule"])
        self._rect(s, bx, by + bh / 2, bw, Emu(9525), fill=self.C["rule"])
        self._txt(s, bx, by + bh + Inches(0.14), bw, Inches(0.3),
                  [(x_label, T_SMALL, True, INK)], align=PP_ALIGN.CENTER)
        self._txt(s, M, by + bh / 2 - Inches(0.15), LBL - Inches(0.14), Inches(0.3),
                  [(y_label, T_SMALL, True, INK)], align=PP_ALIGN.RIGHT)
        # quadrants = (좌상, 우상, 좌하, 우하) — 비어 있는 사분면이 무엇을 뜻하는지
        # 그림이 말하게 한다. 타이틀만으로는 청중이 해석할 수 없다.
        if quadrants:
            qpos = [(bx, by), (bx + bw / 2, by),
                    (bx, by + bh / 2), (bx + bw / 2, by + bh / 2)]
            for q, (qx, qy) in zip(quadrants, qpos):
                if q:
                    self._txt(s, qx + Inches(0.18), qy + Inches(0.14),
                              bw / 2 - Inches(0.36), Inches(0.3),
                              [(q, T_META, True, MUTE)])

        lw, lh = Inches(2.0), Inches(0.3)
        placed = []
        for it in items:
            name, xv, yv = it[0], float(it[1]), float(it[2])
            hot = it[3] if len(it) > 3 else False
            cx = bx + bw * max(0.0, min(1.0, xv))
            cy = by + bh * (1 - max(0.0, min(1.0, yv)))
            d = Inches(0.22) if hot else Inches(0.16)
            self._oval(s, cx - d / 2, cy - d / 2, d, d,
                       self.C["key"] if hot else self.C["sub"])

            # 라벨은 점 아래 중앙. 좌우는 박스 안으로 clamp.
            lx = min(max(cx - lw / 2, bx), bx + bw - lw)
            ly = cy + d / 2 + Inches(0.06)
            # 라벨 박스는 2in 폭에 글자가 가운데 정렬이라, 박스 좌표가 조금 떨어져
            # 있어도 글자끼리는 겹친다. 세로 임계를 넉넉히 잡아야 실제 겹침이 잡힌다.
            for px, py in sorted(placed, key=lambda t: t[1]):
                if abs(px - lx) < lw * 0.9 and abs(py - ly) < lh * 2:
                    ly = py + lh * 1.1
            if ly + lh > by + bh:                      # 바닥을 넘으면 점 위로
                ly = cy - d / 2 - lh - Inches(0.04)
            placed.append((lx, ly))
            # 겹침 회피로 밀려난 라벨은 점과의 소유 관계가 끊긴다. 리더선으로 잇는다.
            if ly > cy + d / 2 + Inches(0.16):
                self._rect(s, cx - Emu(4763), cy + d / 2, Emu(9525),
                           ly - cy - d / 2, fill=self.C["rule"])
            self._txt(s, lx, ly, lw, lh,
                      [(name, T_SMALL, hot, INK if hot else MUTE)], align=PP_ALIGN.CENTER)
        return s

    def shots(self, title, images, captions=None, eyebrow=None, lead=None):
        """스크린샷 그리드. images = [경로, ...] 최대 5."""
        s = self._blank()
        y0 = self._head(s, title, eyebrow, lead)
        imgs = [p for p in images if os.path.exists(p)][:5]
        if not imgs:
            self._txt(s, M, y0, CW, Inches(0.5),
                      [("(이미지 없음)", T_BODY, False, MUTE)])
            return s
        n = len(imgs); g = Inches(0.3)
        cw = (CW - g * (n - 1)) / n
        avail_h = BOTTOM - y0 - Inches(0.4)
        for i, path in enumerate(imgs):
            x = M + (cw + g) * i
            pic = s.shapes.add_picture(path, int(x), int(y0), width=int(cw))
            if pic.height > avail_h:
                ratio = avail_h / pic.height
                pic.height = int(avail_h); pic.width = int(pic.width * ratio)
                pic.left = int(x + (cw - pic.width) / 2)
            # 흰 배경 스크린샷은 흰 슬라이드에 얹으면 경계가 사라진다.
            # 이미지 뒤에 틴트 배킹을 깔아 면을 만든다(§4 목업 예외).
            pad = Inches(0.1)
            bg = self._rect(s, pic.left - pad, pic.top - pad,
                            pic.width + pad * 2, pic.height + pad * 2, fill=self.C["tint"])
            s.shapes._spTree.remove(bg._element)
            s.shapes._spTree.insert(2, bg._element)      # 이미지 뒤로 보낸다
            if captions and i < len(captions):
                # 고정 y가 아니라 이미지 실제 하단 기준. 가로로 긴 이미지에서
                # 캡션이 뚝 떨어져 뜨던 버그.
                self._txt(s, pic.left, pic.top + pic.height + Inches(0.22),
                          pic.width, Inches(0.3),
                          [(captions[i], T_META, False, MUTE)], align=PP_ALIGN.CENTER)
        return s

    def quote(self, text, attribution, eyebrow=None, sub=""):
        """인용·페르소나. attribution은 출처(이름·나이·역할).

        bullets에 넣으면 목소리가 데이터가 된다. statement는 저자(우리)의 주장이고,
        인용은 **출처 귀속이 형태의 본질**이다. 귀속 없는 인용은 인용이 아니다.
        """
        s = self._blank()
        self._head(s, "", eyebrow)
        top = Inches(1.9)
        self._txt(s, M, top, Inches(1.1), Inches(1.2),
                  [("“", 88, True, self.C["pale"])], ls=1.0)
        self._txt(s, M + Inches(0.95), top + Inches(0.2), CW - Inches(1.6), Inches(2.0),
                  [(text, 24, False, INK)], ls=1.5)
        ay = top + Inches(2.5)
        self._rect(s, M + Inches(0.95), ay, Inches(0.6), Emu(19050), fill=self.C["key"])
        self._txt(s, M + Inches(0.95), ay + Inches(0.18), Inches(7), Inches(0.34),
                  [(attribution, T_SMALL, True, self.C["key"])])
        if sub:
            self._txt(s, M + Inches(0.95), ay + Inches(0.58), Inches(9), Inches(0.4),
                      [(sub, T_SMALL, False, MUTE)])
        return s

    def compare(self, title, left, right, rows, eyebrow=None, lead=None):
        """대립 2패널. left/right = 패널 제목, rows = [(축, 좌, 우), ...]

        table은 대립을 중립적 나열로 만들어 어느 쪽이 우리인지 안 보인다.
        cards는 병렬이지 대립이 아니다. 대립은 **가운데 구분 + 한쪽만 강조**로
        성립한다. 우측이 강조 패널(신안·우리)이다.
        """
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        axis_w = CW * 0.18
        pw = (CW - axis_w - Inches(0.3)) / 2
        lx = M + axis_w
        rx = lx + pw + Inches(0.3)
        h = BOTTOM - y
        self._rect(s, lx, y, pw, h, fill=RGBColor(0xF5, 0xF5, 0xF7))
        self._rect(s, rx, y, pw, h, fill=self.C["tint"])
        self._rect(s, rx, y, pw, Inches(0.055), fill=self.C["key"])
        self._txt(s, lx + Inches(0.24), y + Inches(0.24), pw - Inches(0.48), Inches(0.4),
                  [(left, T_BODY, True, MUTE)])
        self._txt(s, rx + Inches(0.24), y + Inches(0.24), pw - Inches(0.48), Inches(0.4),
                  [(right, T_BODY, True, self.C["deep"])])
        n = max(1, len(rows))
        rh = (h - Inches(0.85)) / n
        for i, (axis, lv, rv) in enumerate(rows):
            ry = y + Inches(0.85) + rh * i
            self._txt(s, M, ry + Inches(0.06), axis_w - Inches(0.18), rh,
                      [(axis, T_SMALL, True, MUTE)], align=PP_ALIGN.RIGHT)
            self._txt(s, lx + Inches(0.24), ry + Inches(0.06), pw - Inches(0.48), rh,
                      [(lv, T_SMALL, False, MUTE)], ls=1.4)
            self._txt(s, rx + Inches(0.24), ry + Inches(0.06), pw - Inches(0.48), rh,
                      [(rv, T_SMALL, False, INK)], ls=1.4)
        return s

    def stat(self, number, definition, notes=None, eyebrow=None):
        """큰 숫자 하나. notes = [보조 설명, ...] **3개 초과 시 예외**.

        cards 4장에 넣으면 넷이 동등해져서 북극성이 북극성이 아니게 된다.
        숫자를 문장에 묻는 것도 금지(§2 판정 절차 2번).
        """
        s = self._blank()
        self._head(s, "", eyebrow)
        left_w = CW * 0.62
        size = 108
        while size > 44 and measure_pt(number, size, True) > (left_w / Emu(12700)):
            size -= 4
        top = Inches(2.3)
        self._txt(s, M, top, left_w, Inches(1.8),
                  [(str(number), size, True, self.C["deep"])], ls=1.05)
        self._rect(s, M, top + Inches(1.85), Inches(0.9), Inches(0.05), fill=self.C["key"])
        self._txt(s, M, top + Inches(2.1), left_w, Inches(0.6),
                  [(definition, T_LEAD, False, INK)])
        if notes and len(notes) > 3:
            # 조용히 자르면 만든 사람만 모른다. tree와 같은 원칙으로 멈춘다.
            raise ValueError(f"stat notes는 3개까지다 ({len(notes)}개 받음). "
                             "넘치면 슬라이드를 쪼개거나 문장을 합쳐라.")
        if notes:
            ny = top + Inches(0.1)
            for nt in notes:
                self._txt(s, M + left_w + Inches(0.4), ny, CW - left_w - Inches(0.4),
                          Inches(0.7), [(nt, T_SMALL, False, MUTE)], ls=1.45)
                ny += Inches(0.78)
        return s

    def gate(self, title, items, eyebrow=None, lead=None):
        """체크리스트·게이트. items = [(항목, 상태, 기준), ...]

        상태 = "pass" | "wait" | "fail". 시맨틱 색을 쓰며 브랜드 색으로 덮지 않는다.
        bullets는 상태를 못 담고, table 3열은 "전부 통과해야 한다"는 AND 의미를
        못 보여준다.
        """
        s = self._blank()
        y = self._head(s, title, eyebrow, lead)
        PASS = RGBColor(0x2E, 0x9E, 0x5B)
        WAIT = RGBColor(0xB8, 0x7A, 0x1E)
        FAIL = RGBColor(0xC0, 0x39, 0x2B)
        MARK = {"pass": ("●", PASS), "wait": ("○", WAIT), "fail": ("✕", FAIL)}
        n = max(1, len(items))
        row = min(Inches(0.86), (BOTTOM - y - Inches(0.6)) / n)
        npass = sum(1 for it in items if it[1] == "pass")
        for i, it in enumerate(items):
            name, st = it[0], it[1]
            crit = it[2] if len(it) > 2 else ""
            ry = y + row * i
            glyph, col = MARK.get(st, MARK["wait"])
            self._txt(s, M, ry + Inches(0.04), Inches(0.4), Inches(0.4),
                      [(glyph, T_BODY, True, col)])
            self._txt(s, M + Inches(0.5), ry, CW * 0.42, Inches(0.44),
                      [(name, T_BODY, st == "fail", INK)])
            if crit:
                self._txt(s, M + CW * 0.5, ry + Inches(0.04), CW * 0.5, Inches(0.44),
                          [(crit, T_SMALL, False, MUTE)])
        self._rect(s, M, BOTTOM - Inches(0.5), CW, Emu(9525), fill=self.C["rule"])
        self._txt(s, M, BOTTOM - Inches(0.36), CW, Inches(0.34),
                  [(f"{npass} / {len(items)} 통과 · 전 항목이 통과해야 게이트가 열린다",
                    T_SMALL, True, self.C["key"] if npass == len(items) else MUTE)])
        return s

    def source(self, text, url=None, slide=None):
        """직전 슬라이드에 "정본 → …" 한 줄을 단다. url을 주면 클릭 가능.

        회의에서 본 내용을 나중에 어디서 찾는지 알려주는 줄이다. 이게 없으면
        슬라이드는 읽고 끝나고 아무도 원문을 찾아가지 않는다.

            d.table(...)
            d.source("REQ §10.2 추적 매트릭스", GH + "docs/current/REQ.md")
        """
        s = slide or getattr(self, "_last", None)
        if s is None:
            return None
        run = ("정본 → " + text, T_META, False, self.C["key"] if url else MUTE)
        if url:
            run = run + (url,)
        # 푸터 메타와 **같은 줄**, 쪽번호 앞에 우측 정렬한다.
        # 종전 y(SH-1.12in)는 구분선(SH-0.92in)과 겹쳐, 링크 글자에 취소선처럼
        # 가로선이 그어지고 밑줄까지 더해 한 줄에 선이 3개 겹쳤다. 게다가
        # 콘텐츠 하한(BOTTOM)과 2pt밖에 안 떨어져 표 마지막 행을 눌렀다.
        # 메타 줄 우측은 원래 비어 있던 공간이라 선을 하나 줄이면서 자리도 는다.
        self._txt(s, M + CW * 0.40, SH - Inches(0.78), CW * 0.52, Inches(0.3),
                  [run], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        return s

    def save(self, path):
        self.prs.save(path)
        return path
