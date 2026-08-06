"""
기존 pptx의 레이아웃 답답함을 기계적으로 푼다.

    python3 scripts/deck_polish.py <입력.pptx> [출력.pptx]

`deck_skin.py`(색 치환)와 목적이 다르다. 이건 **여백·위계·구분 방식**을 손본다.
자료형 덱이 답답해 보이는 원인은 대개 색이 아니라 아래 셋이다.

1. **모든 박스에 테두리가 있다** → 격자 감옥. Apple·Toss 계열은 선이 아니라
   **면(연한 채우기)과 여백**으로 구분한다. 채우기가 있는 도형의 선은 지운다.
2. **텍스트가 테두리에 붙어 있다** → 내부 여백을 준다.
3. **타이포 위계 폭이 좁다** → 18/12/10.5/9.5처럼 붙어 있으면 전부 비슷해 보인다.
   본문은 그대로 두고 **제목급만 키워** 위계를 벌린다.

레이아웃 좌표·박스 크기는 건드리지 않는다. 크기를 바꾸면 밀도 높은 슬라이드가
줄바꿈으로 무너지므로, 넘칠 위험이 없는 변경만 한다.

적용 제외: 공모전 기능설명서(지정 양식 — 변경 시 심사 제외).
"""
import sys
from pptx import Presentation
from pptx.util import Pt, Emu

# 제목급만 상향해 위계를 벌린다. 본문(10.5 이하)은 불변 — 키우면 넘친다.
SIZE_MAP = {
    18.0: 23.0,   # 슬라이드 타이틀
    15.0: 18.0,
    14.0: 16.0,
    13.0: 14.5,
    12.0: 13.0,   # 박스 제목
}

PAD = Emu(91440)        # 0.1in — 텍스트 프레임 내부 여백
PAD_TB = Emu(64008)     # 0.07in


def _polish(sh, st):
    # 1) 채우기가 있는 도형의 테두리 제거
    try:
        has_fill = sh.fill.type == 1
    except Exception:
        has_fill = False
    if has_fill:
        try:
            if sh.line.fill.type != 5:      # 5 = background(이미 없음)
                sh.line.fill.background()
                st["line"] += 1
        except Exception:
            pass

    # 2) 텍스트 내부 여백 + 3) 제목급 크기 상향
    if sh.has_text_frame:
        tf = sh.text_frame
        try:
            if has_fill:
                tf.margin_left = tf.margin_right = PAD
                tf.margin_top = tf.margin_bottom = PAD_TB
                st["pad"] += 1
        except Exception:
            pass
        for para in tf.paragraphs:
            for r in para.runs:
                if r.font.size is None:
                    continue
                new = SIZE_MAP.get(round(r.font.size.pt, 1))
                if new:
                    r.font.size = Pt(new)
                    st["size"] += 1

    if sh.shape_type == 6:
        for c in sh.shapes:
            _polish(c, st)


def polish(src, dst):
    prs = Presentation(src)
    st = {"line": 0, "pad": 0, "size": 0}
    for slide in prs.slides:
        for sh in slide.shapes:
            _polish(sh, st)
    prs.save(dst)
    return st


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else src.replace(".pptx", "_polished.pptx")
    st = polish(src, dst)
    print(f"{dst}\n  테두리 제거 {st['line']} · 내부 여백 {st['pad']} · 제목 크기 상향 {st['size']}")
